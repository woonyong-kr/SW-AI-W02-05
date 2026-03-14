"""
BST / 자가균형 트리 / B-Tree 슬라이드 추가 스크립트
기존 Tree 슬라이드 스타일 그대로 적용:
  - 노드박스: 채우기 없음, 검정 테두리, 0.55" x 0.28"
  - 연결선: #6D92AB, 2pt
  - 타이틀: Bold 14pt, 직사각형 박스
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

IN = 914400  # 1 inch in EMU

SRC = '/Users/woonyong/workspace/Krafton-Jungle/SW-AI-W02-05/week3/data-structure/data-structure-basic.pptx'
DST = '/Users/woonyong/workspace/Krafton-Jungle/SW-AI-W02-05/week3/data-structure/data-structure-basic.pptx'

# 공통 치수
NW = Inches(0.55)   # 노드 너비
NH = Inches(0.28)   # 노드 높이
LINE_COLOR = RGBColor(0x6D, 0x92, 0xAB)
LINE_W = 25400      # 2pt
BLACK = RGBColor(0, 0, 0)

def add_node(slide, x_in, y_in, text, fill_rgb=None, font_bold=False, font_size=11, border_color=None):
    """노드 박스 추가. x,y는 인치 단위 중심점"""
    left = Inches(x_in) - NW // 2
    top  = Inches(y_in) - NH // 2
    shape = slide.shapes.add_shape(1, left, top, NW, NH)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.color.rgb = border_color if border_color else BLACK
    shape.line.width = 19050  # 1.5pt
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = font_bold
    run.font.color.rgb = BLACK
    # 수직 중앙 정렬
    tf.margin_top = 0
    tf.margin_bottom = 0
    return shape

def add_line(slide, x1, y1, x2, y2, color=LINE_COLOR, width=LINE_W):
    """두 점을 잇는 선 추가. x,y는 인치 단위"""
    from pptx.util import Inches
    import math
    left   = Inches(min(x1, x2))
    top    = Inches(min(y1, y2))
    right  = Inches(max(x1, x2))
    bottom = Inches(max(y1, y2))
    w = right - left
    h = bottom - top
    if w == 0: w = 1
    if h == 0: h = 1
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_title_box(slide, text, cx=6.67, cy=0.7):
    """타이틀 텍스트 박스"""
    w = Inches(4.0)
    h = Inches(0.42)
    left = Inches(cx) - w // 2
    top  = Inches(cy) - h // 2
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.background()
    shape.line.color.rgb = BLACK
    shape.line.width = 19050
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BLACK

def add_label(slide, x_in, y_in, text, color=RGBColor(0xF9, 0x66, 0x5E), font_size=10):
    """레이블 텍스트 박스"""
    w = Inches(2.5)
    h = Inches(0.3)
    txBox = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), w, h)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color

prs = Presentation(SRC)

# ─────────────────────────────────────────
# 슬라이드 1: 이진 탐색 트리 (BST)
# ─────────────────────────────────────────
blank_layout = prs.slide_layouts[6]  # 빈 슬라이드
slide1 = prs.slides.add_slide(blank_layout)
add_title_box(slide1, '이진 탐색 트리 (Binary Search Tree, BST)', cx=6.67, cy=0.7)

# 노드 위치 정의 (중심 x, y) 인치 단위
# 레벨 0: root=8
# 레벨 1: 3, 10
# 레벨 2: 1, 6, -, 14
# 레벨 3: -, 4, 7, -, 13

# y 위치
Y0, Y1, Y2, Y3 = 1.6, 2.6, 3.6, 4.6

# x 위치 - 전체 너비 13.33" 기준으로 중앙 배치
CX = 6.67
SEP1 = 2.2   # 레벨1 간격
SEP2 = 1.1   # 레벨2 간격
SEP3 = 0.7   # 레벨3 간격

nodes = {
    'root': (CX,           Y0, '8'),
    'L':    (CX - SEP1,    Y1, '3'),
    'R':    (CX + SEP1,    Y1, '10'),
    'LL':   (CX - SEP1 - SEP2, Y2, '1'),
    'LR':   (CX - SEP1 + SEP2, Y2, '6'),
    'RR':   (CX + SEP1 + SEP2, Y2, '14'),
    'LRL':  (CX - SEP1 + SEP2 - SEP3, Y3, '4'),
    'LRR':  (CX - SEP1 + SEP2 + SEP3, Y3, '7'),
    'RRL':  (CX + SEP1 + SEP2 - SEP3, Y3, '13'),
}

for key, (x, y, label) in nodes.items():
    add_node(slide1, x, y, label, font_bold=True, font_size=12)

# 연결선
edges = [
    ('root','L'), ('root','R'),
    ('L','LL'), ('L','LR'),
    ('R','RR'),
    ('LR','LRL'), ('LR','LRR'),
    ('RR','RRL'),
]
for a, b in edges:
    x1,y1,_ = nodes[a]
    x2,y2,_ = nodes[b]
    add_line(slide1, x1, y1+0.14, x2, y2-0.14)

# 규칙 레이블
add_label(slide1, 1.2, 1.4, '← 왼쪽 자식 < 부모', color=RGBColor(0x6D, 0x92, 0xAB), font_size=11)
add_label(slide1, 8.5, 1.4, '오른쪽 자식 > 부모 →', color=RGBColor(0x6D, 0x92, 0xAB), font_size=11)

# 설명 레이블
add_label(slide1, 1.0, 5.3,
    '탐색: 찾는 값이 현재 노드보다 작으면 왼쪽, 크면 오른쪽으로 이동 → O(log n)',
    color=BLACK, font_size=11)
add_label(slide1, 1.0, 5.7,
    '단점: 정렬된 순서로 삽입 시 한쪽으로 치우쳐 O(n)으로 저하',
    color=RGBColor(0xF9, 0x66, 0x5E), font_size=11)

print('Slide 1 (BST) done')

# ─────────────────────────────────────────
# 슬라이드 2: 자가균형 트리 (AVL / Red-Black)
# ─────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_title_box(slide2, '자가균형 트리 (Self-Balancing Tree)', cx=6.67, cy=0.7)

# ── AVL 트리 (왼쪽) ──
avl_cx = 3.5
add_label(slide2, avl_cx - 1.0, 1.1, 'AVL 트리', color=RGBColor(0x6D, 0x92, 0xAB), font_size=12)
add_label(slide2, 0.5, 1.45, '높이 차이 최대 1 유지 | 회전(Rotation)으로 균형 복원', color=BLACK, font_size=10)

avl_nodes = {
    'a4': (avl_cx,        1.9, '4'),
    'a2': (avl_cx - 1.1,  2.8, '2'),
    'a6': (avl_cx + 1.1,  2.8, '6'),
    'a1': (avl_cx - 1.6,  3.7, '1'),
    'a3': (avl_cx - 0.6,  3.7, '3'),
    'a5': (avl_cx + 0.6,  3.7, '5'),
    'a7': (avl_cx + 1.6,  3.7, '7'),
}
for k, (x, y, lbl) in avl_nodes.items():
    add_node(slide2, x, y, lbl, font_bold=True, font_size=12)

avl_edges = [
    ('a4','a2'),('a4','a6'),
    ('a2','a1'),('a2','a3'),
    ('a6','a5'),('a6','a7'),
]
for a, b in avl_edges:
    x1,y1,_ = avl_nodes[a]; x2,y2,_ = avl_nodes[b]
    add_line(slide2, x1, y1+0.14, x2, y2-0.14)

add_label(slide2, 0.5, 4.2, '모든 노드의 left/right 높이 차 ≤ 1', color=RGBColor(0x6D, 0x92, 0xAB), font_size=10)

# ── 레드-블랙 트리 (오른쪽) ──
rb_cx = 9.5
add_label(slide2, rb_cx - 1.3, 1.1, '레드-블랙 트리', color=RGBColor(0xF9, 0x66, 0x5E), font_size=12)
add_label(slide2, 7.0, 1.45, '빨강/검정 규칙으로 균형 유지 | 삽입·삭제 빠름', color=BLACK, font_size=10)

# 레드 노드: #F9665E, 블랙 노드: 채우기 없음+굵은테두리
rb_nodes = {
    'rb7':  (rb_cx,        1.9,  '7',  None),
    'rb3':  (rb_cx - 1.1,  2.8,  '3',  RGBColor(0xF9, 0x66, 0x5E)),
    'rb11': (rb_cx + 1.1,  2.8,  '11', None),
    'rb1':  (rb_cx - 1.6,  3.7,  '1',  None),
    'rb5':  (rb_cx - 0.6,  3.7,  '5',  RGBColor(0xF9, 0x66, 0x5E)),
    'rb9':  (rb_cx + 0.6,  3.7,  '9',  RGBColor(0xF9, 0x66, 0x5E)),
    'rb13': (rb_cx + 1.6,  3.7,  '13', RGBColor(0xF9, 0x66, 0x5E)),
}
for k, (x, y, lbl, fill) in rb_nodes.items():
    fc = RGBColor(0xFF, 0xFF, 0xFF) if fill else None
    add_node(slide2, x, y, lbl, fill_rgb=fill, font_bold=True, font_size=12,
             border_color=RGBColor(0x22,0x22,0x22))
    # 흰 텍스트로 덮어씌우기 (레드 배경일 때)
    if fill:
        # 마지막 추가된 shape의 텍스트 색 변경
        shapes = slide2.shapes
        last = shapes[-1]
        last.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

rb_edges = [
    ('rb7','rb3'),('rb7','rb11'),
    ('rb3','rb1'),('rb3','rb5'),
    ('rb11','rb9'),('rb11','rb13'),
]
for a, b in rb_edges:
    x1,y1,lbl,_ = rb_nodes[a]; x2,y2,lbl2,_ = rb_nodes[b]
    add_line(slide2, x1, y1+0.14, x2, y2-0.14)

add_label(slide2, 7.0, 4.2, '검정=균형 기준 노드 | 빨강=삽입 노드 | 회전 최소화', color=RGBColor(0xF9, 0x66, 0x5E), font_size=10)
add_label(slide2, 1.5, 5.1, '실무 활용: Java TreeMap · C++ STL map · Linux 커널 스케줄러', color=BLACK, font_size=11)

# 구분선
from pptx.util import Inches as I
ln = slide2.shapes.add_connector(1, I(6.5), I(1.2), I(6.5), I(4.8))
ln.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
ln.line.width = 12700

print('Slide 2 (Self-Balancing) done')
